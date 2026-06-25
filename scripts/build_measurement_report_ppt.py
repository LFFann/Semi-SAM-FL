from __future__ import annotations

import argparse
import html
import re
import tempfile
import zipfile
from dataclasses import dataclass
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DEFAULT_DATA_ROOT = ROOT / "Results" / "data_260513"
DEFAULT_TEMPLATE = ROOT / "Measurement" / "大脑外侧裂测量示意图.pptx"
DEFAULT_OUTPUT = ROOT / "Measurement" / "大脑外侧裂测量结果汇报_自动生成.pptx"


SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)

INK = RGBColor(30, 41, 59)
MUTED = RGBColor(100, 116, 139)
BLUE = RGBColor(37, 99, 235)
CYAN = RGBColor(8, 145, 178)
GREEN = RGBColor(22, 163, 74)
AMBER = RGBColor(217, 119, 6)
RED = RGBColor(220, 38, 38)
PAPER = RGBColor(248, 250, 252)
PANEL = RGBColor(255, 255, 255)
LINE = RGBColor(203, 213, 225)


@dataclass
class MetricConfig:
    key: str
    title: str
    definition: str
    value_col: str
    status_col: str
    image_col: str
    unit: str
    template_slide: int
    accent: RGBColor


METRICS: list[MetricConfig] = [
    MetricConfig(
        key="lateral_max_depth",
        title="外侧裂最大深度",
        definition="外侧裂沟底至脑表面最远垂直距离",
        value_col="fissure_depth_px",
        status_col="fissure_depth_measurement_status",
        image_col="lateral_max_depth_image_overlay_path",
        unit="px",
        template_slide=3,
        accent=BLUE,
    ),
    MetricConfig(
        key="lateral_opening_width",
        title="外侧裂开口宽度",
        definition="外侧裂两侧臂开口之间最大宽度",
        value_col="fissure_opening_width_smooth_px",
        status_col="fissure_opening_width_measurement_status",
        image_col="lateral_opening_width_image_overlay_path",
        unit="px",
        template_slide=4,
        accent=CYAN,
    ),
    MetricConfig(
        key="lateral_curvature",
        title="裂隙弯曲度",
        definition="外侧裂走行弧度弯曲程度",
        value_col="fissure_curvature_ratio",
        status_col="fissure_curvature_measurement_status",
        image_col="lateral_curvature_image_overlay_path",
        unit="ratio",
        template_slide=5,
        accent=AMBER,
    ),
    MetricConfig(
        key="lateral_angle",
        title="角度",
        definition="经第三脑室顶点作水平线，外侧裂上臂切线与之夹角（双侧）",
        value_col="fissure_mean_angle_deg",
        status_col="fissure_angle_measurement_status",
        image_col="lateral_angle_image_overlay_path",
        unit="deg",
        template_slide=6,
        accent=GREEN,
    ),
    MetricConfig(
        key="longitudinal_branch_max_depth",
        title="纵裂分支最大深度",
        definition="分支顶点距分支起点的距离",
        value_col="longitudinal_fissure_branch_depth_px",
        status_col="longitudinal_fissure_branch_measurement_status",
        image_col="longitudinal_branch_max_depth_image_overlay_path",
        unit="px",
        template_slide=7,
        accent=RED,
    ),
    MetricConfig(
        key="longitudinal_full_length",
        title="纵裂全长",
        definition="纵裂从头端至尾端整体走行长度（黄线）",
        value_col="longitudinal_fissure_full_length_px",
        status_col="longitudinal_fissure_measurement_status",
        image_col="longitudinal_full_length_image_overlay_path",
        unit="px",
        template_slide=8,
        accent=BLUE,
    ),
    MetricConfig(
        key="longitudinal_area",
        title="纵裂面积",
        definition="AI 自动勾勒纵裂整体区域像素面积（蓝色）",
        value_col="longitudinal_fissure_area_px",
        status_col="longitudinal_fissure_measurement_status",
        image_col="longitudinal_area_image_overlay_path",
        unit="px²",
        template_slide=9,
        accent=CYAN,
    ),
]


def rgb_hex(color: RGBColor) -> str:
    return f"#{color[0]:02x}{color[1]:02x}{color[2]:02x}"


def find_results_dir(data_root: Path) -> Path:
    direct = data_root / "测量结果_PPT全部指标"
    if (direct / "tables" / "measurement_results.csv").exists():
        return direct
    for child in data_root.iterdir():
        if child.is_dir() and (child / "tables" / "measurement_results.csv").exists():
            return child
    raise FileNotFoundError(f"未找到 measurement_results.csv: {data_root}")


def read_template_slide_text(template_pptx: Path) -> dict[int, str]:
    if not template_pptx.exists():
        return {}
    text_by_slide: dict[int, str] = {}
    with zipfile.ZipFile(template_pptx) as zf:
        slide_names = [
            name
            for name in zf.namelist()
            if name.startswith("ppt/slides/slide") and name.endswith(".xml")
        ]
        slide_names.sort(key=lambda x: int(re.search(r"slide(\d+)\.xml", x).group(1)))
        for name in slide_names:
            slide_no = int(re.search(r"slide(\d+)\.xml", name).group(1))
            xml = zf.read(name).decode("utf-8", errors="ignore")
            parts = [html.unescape(t) for t in re.findall(r"<a:t>(.*?)</a:t>", xml)]
            text_by_slide[slide_no] = "\n".join(part.strip() for part in parts if part.strip())
    return text_by_slide


def apply_template_definitions(metrics: list[MetricConfig], template_pptx: Path) -> list[MetricConfig]:
    slide_text = read_template_slide_text(template_pptx)
    updated: list[MetricConfig] = []
    for metric in metrics:
        text = slide_text.get(metric.template_slide, "")
        if not text:
            updated.append(metric)
            continue
        clean = re.sub(r"\n+", "\n", text).strip()
        title = metric.title
        definition = metric.definition
        if "：" in clean:
            title_part, definition_part = clean.split("：", 1)
            title = title_part.strip() or title
            definition = " ".join(definition_part.split()) or definition
        elif ":" in clean:
            title_part, definition_part = clean.split(":", 1)
            title = title_part.strip() or title
            definition = " ".join(definition_part.split()) or definition
        updated.append(
            MetricConfig(
                key=metric.key,
                title=title,
                definition=definition,
                value_col=metric.value_col,
                status_col=metric.status_col,
                image_col=metric.image_col,
                unit=metric.unit,
                template_slide=metric.template_slide,
                accent=metric.accent,
            )
        )
    return updated


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: int = 18,
    color: RGBColor = INK,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    font: str = "Microsoft YaHei",
    valign=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_panel(slide, x: float, y: float, w: float, h: float, fill: RGBColor = PANEL):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = LINE
    shape.line.width = Pt(0.8)
    return shape


def add_rule(slide, x: float, y: float, w: float, color: RGBColor = LINE):
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Pt(1))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line


def add_header(slide, kicker: str, title: str, subtitle: str | None = None, accent: RGBColor = BLUE):
    add_text(slide, kicker.upper(), 0.55, 0.25, 2.5, 0.22, size=8, color=accent, bold=True)
    add_text(slide, title, 0.55, 0.48, 8.5, 0.42, size=24, color=INK, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.58, 0.93, 8.8, 0.35, size=11, color=MUTED)
    add_rule(slide, 0.55, 1.28, 12.2, accent)


def add_footer(slide, page_no: int, source: str):
    add_rule(slide, 0.55, 7.02, 12.2, LINE)
    add_text(slide, source, 0.58, 7.12, 9.7, 0.18, size=7, color=MUTED)
    add_text(slide, f"{page_no:02d}", 12.1, 7.08, 0.6, 0.22, size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def add_image_fit(slide, path: Path, x: float, y: float, w: float, h: float):
    with Image.open(path) as img:
        iw, ih = img.size
    scale = min(w / iw, h / ih)
    draw_w = iw * scale
    draw_h = ih * scale
    left = x + (w - draw_w) / 2
    top = y + (h - draw_h) / 2
    with path.open("rb") as image_stream:
        return slide.shapes.add_picture(image_stream, Inches(left), Inches(top), Inches(draw_w), Inches(draw_h))


def numeric_series(df: pd.DataFrame, metric: MetricConfig) -> pd.Series:
    values = pd.to_numeric(df[metric.value_col], errors="coerce")
    if metric.status_col in df.columns:
        ok = df[metric.status_col].astype(str).str.lower().eq("ok")
        values = values[ok]
    return values.dropna()


def metric_stats(df: pd.DataFrame, metric: MetricConfig) -> dict[str, float | int]:
    values = numeric_series(df, metric)
    total = len(df)
    success = len(values)
    stats: dict[str, float | int] = {
        "total": total,
        "success": success,
        "success_rate": success / total * 100 if total else 0,
    }
    if len(values):
        stats.update(
            mean=float(values.mean()),
            median=float(values.median()),
            std=float(values.std(ddof=0)),
            minimum=float(values.min()),
            maximum=float(values.max()),
        )
    else:
        stats.update(mean=0.0, median=0.0, std=0.0, minimum=0.0, maximum=0.0)
    return stats


def existing_path(value) -> Path | None:
    if pd.isna(value):
        return None
    path = Path(str(value))
    return path if path.exists() else None


def representative_row(df: pd.DataFrame, metric: MetricConfig) -> pd.Series:
    values = pd.to_numeric(df[metric.value_col], errors="coerce")
    mask = values.notna()
    if metric.status_col in df.columns:
        mask &= df[metric.status_col].astype(str).str.lower().eq("ok")
    candidates = df.loc[mask].copy()
    if candidates.empty:
        candidates = df.copy()
    candidates["_value"] = pd.to_numeric(candidates[metric.value_col], errors="coerce")
    median = candidates["_value"].dropna().median()
    candidates["_distance"] = (candidates["_value"] - median).abs()
    candidates = candidates.sort_values(["_distance", "index"])
    for _, row in candidates.iterrows():
        if existing_path(row.get(metric.image_col)):
            return row
    for _, row in df.iterrows():
        if existing_path(row.get(metric.image_col)):
            return row
    return df.iloc[0]


def fmt_num(value: float, unit: str) -> str:
    if unit == "ratio":
        return f"{value:.2f}"
    if unit == "deg":
        return f"{value:.1f}°"
    if unit == "px²":
        return f"{value:,.0f} px²"
    return f"{value:.1f} px"


def make_histogram(series: pd.Series, metric: MetricConfig, out_dir: Path) -> Path:
    plt.rcParams["font.sans-serif"] = ["Microsoft YaHei", "SimHei", "Arial Unicode MS", "DejaVu Sans"]
    plt.rcParams["axes.unicode_minus"] = False
    fig, ax = plt.subplots(figsize=(5.0, 2.1), dpi=180)
    ax.hist(series, bins=18, color=rgb_hex(metric.accent), alpha=0.82, edgecolor="white")
    if len(series):
        ax.axvline(series.median(), color="#0f172a", linestyle="--", linewidth=1.2)
        ax.text(series.median(), ax.get_ylim()[1] * 0.9, "中位数", fontsize=8, color="#0f172a")
    ax.set_title(f"{metric.title}分布", fontsize=10, loc="left")
    ax.set_xlabel(metric.unit, fontsize=8)
    ax.set_ylabel("病例数", fontsize=8)
    ax.tick_params(labelsize=7)
    for spine in ["top", "right"]:
        ax.spines[spine].set_visible(False)
    fig.tight_layout()
    out = out_dir / f"{metric.key}_hist.png"
    fig.savefig(out, transparent=False, facecolor="white")
    plt.close(fig)
    return out


def make_summary_chart(df: pd.DataFrame, metrics: list[MetricConfig], out_dir: Path) -> Path:
    names = []
    means = []
    rates = []
    colors = []
    for metric in metrics:
        stats = metric_stats(df, metric)
        names.append(metric.title)
        means.append(float(stats["mean"]))
        rates.append(float(stats["success_rate"]))
        colors.append(rgb_hex(metric.accent))

    plt.rcParams["font.sans-serif"] = ["Microsoft YaHei", "SimHei", "Arial Unicode MS", "DejaVu Sans"]
    fig, axes = plt.subplots(1, 2, figsize=(10.2, 3.0), dpi=180)
    axes[0].barh(names, means, color=colors, alpha=0.86)
    axes[0].set_title("各指标均值", fontsize=11, loc="left")
    axes[0].tick_params(labelsize=8)
    axes[0].invert_yaxis()
    axes[1].barh(names, rates, color=colors, alpha=0.86)
    axes[1].set_title("各指标可测比例", fontsize=11, loc="left")
    axes[1].set_xlim(0, 105)
    axes[1].tick_params(labelsize=8)
    axes[1].invert_yaxis()
    for ax in axes:
        for spine in ["top", "right"]:
            ax.spines[spine].set_visible(False)
        ax.grid(axis="x", color="#e2e8f0", linewidth=0.7)
    fig.tight_layout()
    out = out_dir / "metric_summary.png"
    fig.savefig(out, facecolor="white")
    plt.close(fig)
    return out


def add_stat_card(slide, label: str, value: str, x: float, y: float, w: float, accent: RGBColor):
    add_panel(slide, x, y, w, 0.8)
    add_text(slide, value, x + 0.14, y + 0.14, w - 0.28, 0.28, size=19, color=accent, bold=True)
    add_text(slide, label, x + 0.14, y + 0.48, w - 0.28, 0.18, size=8, color=MUTED)


def add_metric_table(slide, rows: list[tuple[str, str]], x: float, y: float, w: float, h: float):
    table = slide.shapes.add_table(len(rows), 2, Inches(x), Inches(y), Inches(w), Inches(h)).table
    table.columns[0].width = Inches(w * 0.42)
    table.columns[1].width = Inches(w * 0.58)
    for r, (label, value) in enumerate(rows):
        for c, text in enumerate([label, value]):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(241, 245, 249) if c == 0 else PANEL
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            p = cell.text_frame.paragraphs[0]
            p.text = text
            p.font.name = "Microsoft YaHei"
            p.font.size = Pt(8.5)
            p.font.bold = c == 0
            p.font.color.rgb = MUTED if c == 0 else INK


def build_cover(prs: Presentation, df: pd.DataFrame, metrics: list[MetricConfig], source: str, out_dir: Path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    add_text(slide, "大脑外侧裂测量结果汇报", 0.7, 0.75, 8.8, 0.62, size=32, color=INK, bold=True)
    add_text(slide, "依据《大脑外侧裂测量示意图》自动汇总 7 项测量指标", 0.74, 1.42, 8.6, 0.35, size=14, color=MUTED)
    add_rule(slide, 0.72, 1.95, 5.6, BLUE)

    total = len(df)
    lateral_ok = int(df["fissure_measurement_status"].astype(str).str.lower().eq("ok").sum())
    longitudinal_ok = int(df["longitudinal_fissure_measurement_status"].astype(str).str.lower().eq("ok").sum())
    angle_ok = int(df["fissure_angle_measurement_status"].astype(str).str.lower().eq("ok").sum())
    add_stat_card(slide, "总病例数", f"{total}", 0.8, 2.45, 2.0, BLUE)
    add_stat_card(slide, "外侧裂可测", f"{lateral_ok}", 3.05, 2.45, 2.0, CYAN)
    add_stat_card(slide, "角度双侧可测", f"{angle_ok}", 5.3, 2.45, 2.0, GREEN)
    add_stat_card(slide, "纵裂可测", f"{longitudinal_ok}", 7.55, 2.45, 2.0, AMBER)

    chart = make_summary_chart(df, metrics, out_dir)
    add_panel(slide, 0.78, 3.55, 8.8, 2.9)
    add_image_fit(slide, chart, 0.95, 3.78, 8.45, 2.45)

    add_panel(slide, 10.0, 0.72, 2.65, 5.72)
    add_text(slide, "页序继承示意图", 10.25, 1.0, 2.1, 0.26, size=13, color=INK, bold=True)
    lines = ["01 封面", "02 形态总览"] + [f"{i+3:02d} {m.title}" for i, m in enumerate(metrics)]
    add_text(slide, "\n".join(lines), 10.25, 1.45, 2.1, 4.3, size=10, color=MUTED)
    add_footer(slide, 1, source)


def build_morphology_slide(prs: Presentation, df: pd.DataFrame, source: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    add_header(slide, "overview", "外侧裂及大脑纵裂形态", "同一代表病例展示原图、外侧裂测量叠加、纵裂全长与纵裂面积。", BLUE)

    metric = METRICS[1]
    row = representative_row(df, metric)
    image_paths = [
        ("原始图像", existing_path(row.get("source_image_path"))),
        ("外侧裂开口宽度", existing_path(row.get("lateral_opening_width_image_overlay_path"))),
        ("纵裂全长", existing_path(row.get("longitudinal_full_length_image_overlay_path"))),
        ("纵裂面积", existing_path(row.get("longitudinal_area_image_overlay_path"))),
    ]
    positions = [(0.75, 1.65), (3.9, 1.65), (7.05, 1.65), (10.2, 1.65)]
    for (label, path), (x, y) in zip(image_paths, positions):
        add_panel(slide, x, y, 2.65, 3.25)
        if path:
            add_image_fit(slide, path, x + 0.12, y + 0.18, 2.41, 2.65)
        add_text(slide, label, x + 0.12, y + 2.92, 2.41, 0.2, size=9, color=INK, bold=True, align=PP_ALIGN.CENTER)

    add_panel(slide, 0.75, 5.35, 11.95, 0.85)
    add_text(
        slide,
        "自动测量结果按 7 个指标分别保存原图叠加与掩膜叠加图；本汇报优先展示原图叠加图，便于导师核对测量线、角度、面积区域与临床图像背景的对应关系。",
        0.98,
        5.58,
        11.4,
        0.34,
        size=12,
        color=INK,
    )
    add_footer(slide, 2, source)


def build_metric_slide(
    prs: Presentation,
    df: pd.DataFrame,
    metric: MetricConfig,
    page_no: int,
    source: str,
    out_dir: Path,
):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    add_header(slide, "metric result", metric.title, metric.definition, metric.accent)

    stats = metric_stats(df, metric)
    series = numeric_series(df, metric)
    row = representative_row(df, metric)
    image_path = existing_path(row.get(metric.image_col))
    row_value = pd.to_numeric(pd.Series([row.get(metric.value_col)]), errors="coerce").iloc[0]

    add_panel(slide, 0.72, 1.55, 6.15, 4.9)
    if image_path:
        add_image_fit(slide, image_path, 0.9, 1.75, 5.78, 4.22)
    add_text(
        slide,
        f"代表病例 index={int(row.get('index', 0))}，测量值 {fmt_num(float(row_value), metric.unit) if pd.notna(row_value) else 'NA'}",
        0.93,
        6.08,
        5.75,
        0.18,
        size=8,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )

    add_stat_card(slide, "可测 / 总数", f"{int(stats['success'])}/{int(stats['total'])}", 7.18, 1.55, 1.7, metric.accent)
    add_stat_card(slide, "可测比例", f"{stats['success_rate']:.1f}%", 9.05, 1.55, 1.55, metric.accent)
    add_stat_card(slide, "均值", fmt_num(float(stats["mean"]), metric.unit), 10.78, 1.55, 1.7, metric.accent)

    table_rows = [
        ("中位数", fmt_num(float(stats["median"]), metric.unit)),
        ("标准差", fmt_num(float(stats["std"]), metric.unit)),
        ("最小值", fmt_num(float(stats["minimum"]), metric.unit)),
        ("最大值", fmt_num(float(stats["maximum"]), metric.unit)),
    ]
    add_metric_table(slide, table_rows, 7.18, 2.63, 5.3, 1.45)

    if len(series):
        chart = make_histogram(series, metric, out_dir)
        add_panel(slide, 7.18, 4.35, 5.3, 1.72)
        add_image_fit(slide, chart, 7.32, 4.48, 5.02, 1.45)

    add_text(
        slide,
        f"字段：{metric.value_col}；状态字段：{metric.status_col}。统计仅纳入状态为 ok 且数值非空的病例。",
        7.25,
        6.25,
        5.15,
        0.26,
        size=8,
        color=MUTED,
    )
    add_footer(slide, page_no, source)


def build_deck(results_dir: Path, template_pptx: Path, output_pptx: Path):
    csv_path = results_dir / "tables" / "measurement_results.csv"
    if not csv_path.exists():
        raise FileNotFoundError(csv_path)
    with csv_path.open("r", encoding="utf-8-sig", newline="") as handle:
        df = pd.read_csv(handle)
    metrics = apply_template_definitions(METRICS, template_pptx)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    source = f"数据：{csv_path.name}；定义来源：{template_pptx.name if template_pptx.exists() else '内置指标定义'}"
    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    with tempfile.TemporaryDirectory(prefix="measurement_report_") as tmp:
        tmp_dir = Path(tmp)
        build_cover(prs, df, metrics, source, tmp_dir)
        build_morphology_slide(prs, df, source)
        for i, metric in enumerate(metrics, start=3):
            build_metric_slide(prs, df, metric, i, source, tmp_dir)
        with output_pptx.open("wb") as handle:
            prs.save(handle)
    return output_pptx


def parse_args():
    parser = argparse.ArgumentParser(description="自动生成大脑外侧裂/纵裂测量结果汇报 PPT。")
    parser.add_argument(
        "--data-root",
        type=Path,
        default=DEFAULT_DATA_ROOT,
        help="包含测量结果目录的数据根目录，默认 Results/data_260513。",
    )
    parser.add_argument(
        "--results-dir",
        type=Path,
        default=None,
        help="测量结果_PPT全部指标目录；未指定时自动从 data-root 查找 measurement_results.csv。",
    )
    parser.add_argument(
        "--template-pptx",
        type=Path,
        default=DEFAULT_TEMPLATE,
        help="导师给出的测量示意图 PPTX，用于读取指标定义。",
    )
    parser.add_argument(
        "--output",
        type=Path,
        default=DEFAULT_OUTPUT,
        help="输出 PPTX 路径。",
    )
    return parser.parse_args()


def main():
    args = parse_args()
    results_dir = args.results_dir or find_results_dir(args.data_root)
    output = build_deck(results_dir, args.template_pptx, args.output)
    print(f"Generated: {output}")


if __name__ == "__main__":
    main()
